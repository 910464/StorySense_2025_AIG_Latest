#!/usr/bin/env python3
"""
Context Processor for StorySense with AWS Titan Embeddings
Stores embeddings directly in Postgre SQL vector database
"""

import os
import sys
import logging
import time
import json
import argparse
from pathlib import Path
import pandas as pd
from datetime import datetime
import hashlib
from colorama import Fore, Style, init
from dotenv import load_dotenv

# Fix imports by adding parent directory to path if needed
current_dir = Path(__file__).parent.absolute()
parent_dir = current_dir.parent
if str(parent_dir) not in sys.path:
    sys.path.append(str(parent_dir))

# Now try to import the required modules
try:
    # First try direct imports (if running from US_to_StorySense directory)
    from StorySense.PGVectorConnector import PGVectorConnector
    from StorySense.metrics_manager import MetricsManager
    # from env_manager import EnvManager
except ImportError:
    try:
        # Then try package imports (if US_to_StorySense is installed)
        from src.context_handler.context_storage_handler.pgvector_connector import PGVectorConnector
        from src.metrics.metrics_manager import  MetricsManager
        # from src.configuration_handler.env_manager import EnvManager
        from src.llm_layer.image_parser_llm import ImageParserLLM
    except ImportError:
        print(f"{Fore.RED}Error: Cannot import required modules. Please check your installation.{Style.RESET_ALL}")
        print(f"{Fore.YELLOW}Current directory: {current_dir}{Style.RESET_ALL}")
        print(f"{Fore.YELLOW}Python path: {sys.path}{Style.RESET_ALL}")
        sys.exit(1)

# Initialize colorama
init()

# Set up logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# Load environment variables from .env file
load_dotenv()


class EnhancedContextProcessor:
    """Enhanced context processor with PostgreSQL vector storage"""

    def __init__(self, context_folder_path):
        self.context_folder_path = Path(context_folder_path)

        # Initialize environment manager
        # self.env_manager = EnvManager()

        # Initialize metrics manager
        self.metrics_manager = MetricsManager()

        # Initialize image parser
        self.image_parser = ImageParserLLM(metrics_manager=self.metrics_manager)



        # Initialize PGVector connectors for different context types
        self.context_collections = {
            'business_rules': PGVectorConnector(collection_name="business_rules", metrics_manager=self.metrics_manager),
            'requirements': PGVectorConnector(collection_name="requirements", metrics_manager=self.metrics_manager),
            'documentation': PGVectorConnector(collection_name="documentation", metrics_manager=self.metrics_manager),
            'policies': PGVectorConnector(collection_name="policies", metrics_manager=self.metrics_manager),
            'examples': PGVectorConnector(collection_name="examples", metrics_manager=self.metrics_manager),
            'glossary': PGVectorConnector(collection_name="glossary", metrics_manager=self.metrics_manager),
            'wireframes': PGVectorConnector(collection_name="wireframes", metrics_manager=self.metrics_manager)
        }

        # Supported formats
        self.text_formats = ['.xlsx', '.xls', '.csv', '.docx', '.doc', '.pdf', '.json', '.xml', '.txt', '.md', '.pptx',
                             '.ppt']
        self.image_formats = ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp']

        # Create output directory for logs
        self.output_dir = Path('../Output/ProcessedContext')
        self.output_dir.mkdir(exist_ok=True, parents=True)

        # Document registry
        self.registry_file = Path('../Data/DocumentRegistry/document_registry.json')
        os.makedirs('../Data/DocumentRegistry', exist_ok=True)
        self.document_registry = self._load_document_registry()

        # Processing statistics
        self.stats = {
            'total_files': 0,
            'processed_files': 0,
            'skipped_files': 0,
            'failed_files': 0,
            'total_chunks': 0,
            'processing_time': 0,
            'by_category': {}
        }

        # Test database connection
        self._test_db_connection()

    def _test_db_connection(self):
        """Test connection to PostgreSQL database"""
        print(f"{Fore.CYAN}Testing database connection...{Style.RESET_ALL}")
        try:
            # Use the business_rules connector to test the connection
            connector = self.context_collections['business_rules']
            connector.diagnose_database()
            print(f"{Fore.GREEN}✅ Successfully connected to PostgreSQL database{Style.RESET_ALL}")
        except Exception as e:
            print(f"{Fore.RED}❌ Failed to connect to PostgreSQL database: {e}{Style.RESET_ALL}")
            print(f"{Fore.YELLOW}Please check your database configuration in .env file{Style.RESET_ALL}")
            sys.exit(1)

    def _load_document_registry(self):
        """Load the document registry that tracks processed files"""
        if self.registry_file.exists():
            try:
                with open(self.registry_file, 'r', encoding='utf-8') as f:
                    return json.load(f)
            except (json.JSONDecodeError, Exception) as e:
                logger.warning(f"Could not load document registry: {e}. Creating new registry.")
                return {}
        return {}

    def _save_document_registry(self):
        """Save the document registry"""
        try:
            with open(self.registry_file, 'w', encoding='utf-8') as f:
                json.dump(self.document_registry, f, indent=2, default=str)
        except Exception as e:
            logger.error(f"Could not save document registry: {e}")

    def _calculate_file_hash(self, file_path):
        """Calculate MD5 hash of a file for duplicate detection"""
        hash_md5 = hashlib.md5()
        try:
            with open(file_path, "rb") as f:
                for chunk in iter(lambda: f.read(4096), b""):
                    hash_md5.update(chunk)
            return hash_md5.hexdigest()
        except Exception as e:
            logger.error(f"Could not calculate hash for {file_path}: {e}")
            return ""

    def _is_document_processed(self, file_path):
        """Check if a document has already been processed"""
        file_hash = self._calculate_file_hash(file_path)
        file_name = Path(file_path).name
        file_size = os.path.getsize(file_path)
        file_mtime = os.path.getmtime(file_path)

        # Check by hash first (most reliable)
        for doc_id, doc_info in self.document_registry.items():
            if doc_info.get('file_hash') == file_hash and file_hash:
                return {
                    'exists': True,
                    'reason': 'identical_content',
                    'existing_doc_id': doc_id,
                    'existing_info': doc_info
                }

        # Check by name, size, and modification time
        for doc_id, doc_info in self.document_registry.items():
            if (doc_info.get('file_name') == file_name and
                    doc_info.get('file_size') == file_size and
                    doc_info.get('file_mtime') == file_mtime):
                return {
                    'exists': True,
                    'reason': 'same_file_attributes',
                    'existing_doc_id': doc_id,
                    'existing_info': doc_info
                }

        return {'exists': False}

    def _register_document(self, file_path, processing_result):
        """Register a processed document in the registry"""
        doc_id = f"doc_{datetime.now().strftime('%Y%m%d_%H%M%S')}_{len(self.document_registry)}"

        self.document_registry[doc_id] = {
            'file_name': Path(file_path).name,
            'file_path': str(file_path),
            'file_hash': self._calculate_file_hash(file_path),
            'file_size': os.path.getsize(file_path),
            'file_mtime': os.path.getmtime(file_path),
            'file_type': Path(file_path).suffix.lower(),
            'processed_date': datetime.now().isoformat(),
            'document_count': processing_result.get('document_count', 0),
            'collections': processing_result.get('collections', []),
            'context_type': processing_result.get('context_type', 'unknown'),
            'processing_time': processing_result.get('processing_time', 0),
            'success': processing_result.get('success', False)
        }

        self._save_document_registry()
        return doc_id

    def process_all_context_files(self, force_reprocess=False):
        """Process all context files in the folder and subfolders"""
        print(f"\n{Fore.CYAN}🚀 Starting Enhanced Context Processing{Style.RESET_ALL}")
        print(f"📁 Context Folder: {self.context_folder_path}")
        print(f"🔄 Force Reprocess: {'Yes' if force_reprocess else 'No'}")
        print(f"💾 Storage: PostgreSQL Vector Database")

        start_time = time.time()

        if not self.context_folder_path.exists():
            raise FileNotFoundError(f"Context folder not found: {self.context_folder_path}")

        # Discover all files
        all_files = self._discover_files()
        self.stats['total_files'] = sum(len(files) for files in all_files.values())

        print(f"📊 Found {self.stats['total_files']} files to process")

        # Process files by category
        for category, files in all_files.items():
            if files:
                print(f"\n{Fore.YELLOW}📂 Processing {category} files ({len(files)} files){Style.RESET_ALL}")
                self._process_category_files(category, files, force_reprocess)

        # Calculate final statistics
        self.stats['processing_time'] = time.time() - start_time
        self._display_final_stats()

        return self.stats

    def _discover_files(self):
        """Discover and categorize all files in the context folder"""
        categorized_files = {
            'business_rules': [],
            'requirements': [],
            'documentation': [],
            'policies': [],
            'examples': [],
            'glossary': [],
            'wireframes': []
        }

        for root, dirs, files in os.walk(self.context_folder_path):
            category_hint = Path(root).name.lower()

            for file in files:
                file_path = Path(root) / file
                file_ext = file_path.suffix.lower()

                # Skip hidden files and system files
                if file.startswith('.') or file.startswith('~'):
                    continue

                # Skip files that are too large
                max_size_mb = float(os.getenv('MAX_FILE_SIZE_MB', '500'))
                if os.path.getsize(file_path) > max_size_mb * 1024 * 1024:
                    logger.warning(f"Skipping large file: {file_path} (exceeds {max_size_mb}MB)")
                    continue

                # Categorize based on extension and folder name
                category = self._categorize_file(file_path, category_hint)
                categorized_files[category].append(file_path)

        return categorized_files

    def _categorize_file(self, file_path, folder_hint):
        """Categorize file based on extension and folder context"""
        file_ext = file_path.suffix.lower()
        file_name = file_path.name.lower()

        # Check folder name first for strongest hint
        if any(hint in folder_hint for hint in ['business', 'rule', 'logic', 'constraint']):
            return 'business_rules'
        elif any(hint in folder_hint for hint in ['requirement', 'spec', 'functional', 'non-functional']):
            return 'requirements'
        elif any(hint in folder_hint for hint in ['doc', 'guide', 'manual', 'readme']):
            return 'documentation'
        elif any(hint in folder_hint for hint in ['policy', 'procedure', 'guideline', 'standard']):
            return 'policies'
        elif any(hint in folder_hint for hint in ['example', 'template', 'sample', 'demo']):
            return 'examples'
        elif any(hint in folder_hint for hint in ['glossary', 'term', 'definition', 'acronym']):
            return 'glossary'
        elif any(hint in folder_hint for hint in ['wireframe', 'mockup', 'ui', 'design']):
            return 'wireframes'

        # Then check file name
        if any(keyword in file_name for keyword in ['business', 'rule', 'logic', 'constraint']):
            return 'business_rules'
        elif any(keyword in file_name for keyword in ['requirement', 'spec', 'functional', 'non-functional']):
            return 'requirements'
        elif any(keyword in file_name for keyword in ['doc', 'guide', 'manual', 'readme']):
            return 'documentation'
        elif any(keyword in file_name for keyword in ['policy', 'procedure', 'guideline', 'standard']):
            return 'policies'
        elif any(keyword in file_name for keyword in ['example', 'template', 'sample', 'demo']):
            return 'examples'
        elif any(keyword in file_name for keyword in ['glossary', 'term', 'definition', 'acronym']):
            return 'glossary'
        elif any(keyword in file_name for keyword in ['wireframe', 'mockup', 'ui', 'design']):
            return 'wireframes'

        # Default based on file extension
        if file_ext in self.image_formats:
            return 'wireframes' if 'ui' in file_name or 'screen' in file_name else 'documentation'
        elif file_ext in ['.xlsx', '.xls', '.csv']:
            return 'glossary' if 'glossary' in file_name else 'business_rules'
        elif file_ext in ['.pdf', '.docx', '.doc']:
            return 'documentation'
        elif file_ext in ['.pptx', '.ppt']:
            return 'requirements'
        else:
            return 'documentation'  # Default category

    def _process_category_files(self, category, files, force_reprocess):
        """Process files in a specific category"""
        if category not in self.stats['by_category']:
            self.stats['by_category'][category] = {
                'total': len(files),
                'processed': 0,
                'skipped': 0,
                'failed': 0
            }

        image_batch = []
        MAX_BATCH_SIZE_MB = 2.0

        for file_path in files:
            try:
                print(f"  📄 Processing: {file_path.name}")

                if not force_reprocess:
                    duplicate_check = self._is_document_processed(file_path)
                    if duplicate_check['exists']:
                        print(f"    ⏭️  Skipping {Path(file_path).name} - {duplicate_check['reason']}")
                        self.stats['skipped_files'] += 1
                        self.stats['by_category'][category]['skipped'] += 1
                        continue

                file_ext = file_path.suffix.lower()

                if file_ext in self.image_formats:
                    file_size_mb = os.path.getsize(file_path) / (1024 * 1024)
                    current_batch_size_mb = sum(os.path.getsize(fp) / (1024 * 1024) for fp in image_batch)

                    if current_batch_size_mb + file_size_mb > MAX_BATCH_SIZE_MB:
                        if image_batch:
                            self._process_image_batch(image_batch, category)
                            image_batch = []
                    
                    image_batch.append(file_path)

                elif file_ext in ['.txt', '.md']:
                    result = self._process_text_file(file_path, category)
                elif file_ext in ['.pdf']:
                    result = self._process_pdf_file(file_path, category)
                elif file_ext in ['.docx', '.doc']:
                    result = self._process_doc_file(file_path, category)
                elif file_ext in ['.xlsx', '.xls', '.csv']:
                    result = self._process_excel_file(file_path, category)
                elif file_ext in ['.pptx', '.ppt']:
                    result = self._process_presentation_file(file_path, category)
                else:
                    print(f"    ⚠️  Skipping unsupported file type: {file_ext}")
                    self.stats['skipped_files'] += 1
                    self.stats['by_category'][category]['skipped'] += 1
                    continue
                
                if file_ext not in self.image_formats:
                    self._handle_processing_result(result, file_path, category)

            except Exception as e:
                logger.error(f"Error processing {file_path}: {str(e)}")
                print(f"    ❌ Failed: {str(e)}")
                self.stats['failed_files'] += 1
                self.stats['by_category'][category]['failed'] += 1
        
        if image_batch:
            self._process_image_batch(image_batch, category)

    def _process_image_batch(self, image_paths, category):
        """Process a batch of image files."""
        print(f"  🖼️ Processing batch of {len(image_paths)} images...")
        try:
            parsed_texts = self.image_parser.parse_image_batch(image_paths)
            for file_path, parsed_text in zip(image_paths, parsed_texts):
                result = {
                    'success': not parsed_text.startswith("Error:"),
                    'document_count': 0,
                    'collections': [category],
                    'context_type': category,
                    'processing_time': 0, # This would need to be timed per image if needed
                    'error': parsed_text if parsed_text.startswith("Error:") else None
                }

                if result['success']:
                    document = {
                        'text': parsed_text,
                        'metadata': {
                            'source_file': file_path.name,
                            'file_type': 'image',
                            'category': category,
                            'file_size': os.path.getsize(file_path),
                        }
                    }
                    connector = self.context_collections[category]
                    if connector.vector_store_documents([document]):
                        result['document_count'] = 1
                    else:
                        result['success'] = False
                        result['error'] = "Failed to store document in vector database"

                self._handle_processing_result(result, file_path, category)

        except Exception as e:
            logger.error(f"Error processing image batch: {e}")
            for file_path in image_paths:
                self.stats['failed_files'] += 1
                self.stats['by_category'][category]['failed'] += 1
                print(f"    ❌ Failed: {file_path.name} due to batch processing error.")

    def _handle_processing_result(self, result, file_path, category):
        """Handles the result of file processing for stats and registration."""
        if result['success']:
            self._register_document(file_path, result)
            self.stats['processed_files'] += 1
            self.stats['by_category'][category]['processed'] += 1
            self.stats['total_chunks'] += result['document_count']
            print(f"    ✅ Successfully processed {result['document_count']} chunks and stored in PostgreSQL")
        else:
            self.stats['failed_files'] += 1
            self.stats['by_category'][category]['failed'] += 1
            print(f"    ❌ Failed: {result.get('error', 'Unknown error')}")

    def _process_text_file(self, file_path, category):
        """Process text files (.txt, .md) and store in PostgreSQL"""
        start_time = time.time()
        result = {
            'success': False,
            'document_count': 0,
            'collections': [category],
            'context_type': category,
            'processing_time': 0,
            'error': None
        }

        try:
            # Read the file
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
            except UnicodeDecodeError:
                # Try with different encoding
                with open(file_path, 'r', encoding='latin-1') as f:
                    content = f.read()

            # Create document for vector storage
            document = {
                'text': content,
                'metadata': {
                    'source_file': file_path.name,
                    'file_type': 'text',
                    'category': category,
                    'file_size': os.path.getsize(file_path),
                    'character_count': len(content),
                    'line_count': len(content.split('\n'))
                }
            }

            # Store in PostgreSQL vector database
            connector = self.context_collections[category]
            success = connector.vector_store_documents([document])

            if success:
                result['success'] = True
                result['document_count'] = 1
            else:
                result['error'] = "Failed to store document in vector database"

        except Exception as e:
            result['error'] = str(e)
            logger.error(f"Error processing text file {file_path}: {e}")

        finally:
            result['processing_time'] = time.time() - start_time

        return result

    def _process_pdf_file(self, file_path, category):
        """Process PDF files and store in PostgreSQL"""
        start_time = time.time()
        result = {
            'success': False,
            'document_count': 0,
            'collections': [category],
            'context_type': category,
            'processing_time': 0,
            'error': None
        }

        try:
            # Try to extract text from PDF
            try:
                import PyPDF2
                documents = []

                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)

                    # Process each page as a separate document
                    for page_num, page in enumerate(pdf_reader.pages):
                        text = page.extract_text()
                        if text.strip():
                            documents.append({
                                'text': text.strip(),
                                'metadata': {
                                    'source_file': Path(file_path).name,
                                    'file_type': 'pdf',
                                    'category': category,
                                    'page_number': page_num + 1,
                                    'total_pages': len(pdf_reader.pages)
                                }
                            })

                    # If no pages were extracted, create a placeholder document
                    if not documents:
                        documents.append({
                            'text': f"PDF file: {file_path.name} (no text extracted)",
                            'metadata': {
                                'source_file': Path(file_path).name,
                                'file_type': 'pdf',
                                'category': category,
                                'total_pages': len(pdf_reader.pages),
                                'extraction_failed': True
                            }
                        })

            except Exception as e:
                logger.warning(f"Could not extract text from PDF: {e}")
                documents = [{
                    'text': f"PDF file: {file_path.name} (text extraction failed: {str(e)})",
                    'metadata': {
                        'source_file': Path(file_path).name,
                        'file_type': 'pdf',
                        'category': category,
                        'file_size': os.path.getsize(file_path),
                        'extraction_failed': True
                    }
                }]

            # Store in PostgreSQL vector database
            connector = self.context_collections[category]
            success = connector.vector_store_documents(documents)

            if success:
                result['success'] = True
                result['document_count'] = len(documents)
            else:
                result['error'] = "Failed to store documents in vector database"

        except Exception as e:
            result['error'] = str(e)
            logger.error(f"Error processing PDF file {file_path}: {e}")

        finally:
            result['processing_time'] = time.time() - start_time

        return result

    def _process_doc_file(self, file_path, category):
        """Process Word documents and store in PostgreSQL"""
        start_time = time.time()
        result = {
            'success': False,
            'document_count': 0,
            'collections': [category],
            'context_type': category,
            'processing_time': 0,
            'error': None
        }

        try:
            # Try to extract text from Word document
            try:
                import docx
                doc = docx.Document(file_path)

                # Extract paragraphs
                full_text = []
                for paragraph in doc.paragraphs:
                    if paragraph.text.strip():
                        full_text.append(paragraph.text.strip())

                # Extract tables
                table_content = []
                for table in doc.tables:
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            table_content.append(" | ".join(row_text))

                # Combine content
                all_content = "\n".join(full_text)
                if table_content:
                    all_content += "\n\nTables:\n" + "\n".join(table_content)

                if all_content.strip():
                    document = {
                        'text': all_content.strip(),
                        'metadata': {
                            'source_file': Path(file_path).name,
                            'file_type': 'docx',
                            'category': category,
                            'paragraph_count': len(full_text),
                            'table_count': len(doc.tables)
                        }
                    }
                else:
                    document = {
                        'text': f"Word document: {file_path.name} (no text content)",
                        'metadata': {
                            'source_file': Path(file_path).name,
                            'file_type': 'docx',
                            'category': category,
                            'empty_document': True
                        }
                    }

                documents = [document]

            except Exception as e:
                logger.warning(f"Could not extract text from Word document: {e}")
                documents = [{
                    'text': f"Word document: {file_path.name} (text extraction failed: {str(e)})",
                    'metadata': {
                        'source_file': Path(file_path).name,
                        'file_type': 'docx',
                        'category': category,
                        'file_size': os.path.getsize(file_path),
                        'extraction_failed': True
                    }
                }]

            # Store in PostgreSQL vector database
            connector = self.context_collections[category]
            success = connector.vector_store_documents(documents)

            if success:
                result['success'] = True
                result['document_count'] = len(documents)
            else:
                result['error'] = "Failed to store documents in vector database"

        except Exception as e:
            result['error'] = str(e)
            logger.error(f"Error processing Word document {file_path}: {e}")

        finally:
            result['processing_time'] = time.time() - start_time

        return result

    def _process_excel_file(self, file_path, category):
        """Process Excel files and store in PostgreSQL"""
        start_time = time.time()
        result = {
            'success': False,
            'document_count': 0,
            'collections': [category],
            'context_type': category,
            'processing_time': 0,
            'error': None
        }

        try:
            # Try to extract data from Excel
            try:
                df = pd.read_excel(file_path)
                documents = []

                # Process each row as a separate document if it's a structured table
                if len(df) > 0 and len(df.columns) > 1:
                    for idx, row in df.iterrows():
                        content = self._row_to_text(row)
                        if content.strip():
                            documents.append({
                                'text': content,
                                'metadata': {
                                    'source_file': file_path.name,
                                    'file_type': 'excel',
                                    'category': category,
                                    'row_index': idx
                                }
                            })

                # Also add the entire sheet as a single document for context
                text_content = df.to_string(index=False)
                if text_content.strip():
                    documents.append({
                        'text': text_content,
                        'metadata': {
                            'source_file': file_path.name,
                            'file_type': 'excel',
                            'category': category,
                            'row_count': len(df),
                            'column_count': len(df.columns),
                            'full_sheet': True
                        }
                    })

                # If no content was extracted, create a placeholder document
                if not documents:
                    documents.append({
                        'text': f"Excel file: {file_path.name} (no data extracted)",
                        'metadata': {
                            'source_file': file_path.name,
                            'file_type': 'excel',
                            'category': category,
                            'empty_file': True
                        }
                    })

            except Exception as e:
                logger.warning(f"Could not extract data from Excel: {e}")
                documents = [{
                    'text': f"Excel file: {file_path.name} (data extraction failed: {str(e)})",
                    'metadata': {
                        'source_file': file_path.name,
                        'file_type': 'excel',
                        'category': category,
                        'file_size': os.path.getsize(file_path),
                        'extraction_failed': True
                    }
                }]

            # Store in PostgreSQL vector database
            connector = self.context_collections[category]
            success = connector.vector_store_documents(documents)

            if success:
                result['success'] = True
                result['document_count'] = len(documents)
            else:
                result['error'] = "Failed to store documents in vector database"

        except Exception as e:
            result['error'] = str(e)
            logger.error(f"Error processing Excel file {file_path}: {e}")

        finally:
            result['processing_time'] = time.time() - start_time

        return result

    def _row_to_text(self, row):
        """Convert DataFrame row to text content"""
        content_parts = []
        for col, val in row.items():
            if pd.notna(val) and str(val).strip():
                content_parts.append(f"{col}: {val}")
        return "\n".join(content_parts)

    def _process_presentation_file(self, file_path, category):
        """Process PowerPoint presentations and store in PostgreSQL"""
        start_time = time.time()
        result = {
            'success': False,
            'document_count': 0,
            'collections': [category],
            'context_type': category,
            'processing_time': 0,
            'error': None
        }

        try:
            # Try to extract text from PowerPoint
            try:
                from pptx import Presentation
                documents = []
                prs = Presentation(file_path)

                for slide_num, slide in enumerate(prs.slides):
                    slide_text = []

                    # Extract text from shapes
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text.strip())

                    if slide_text:
                        documents.append({
                            'text': "\n".join(slide_text),
                            'metadata': {
                                'source_file': Path(file_path).name,
                                'file_type': 'pptx',
                                'category': category,
                                'slide_number': slide_num + 1,
                                'total_slides': len(prs.slides)
                            }
                        })

                # If no slides were extracted, create a placeholder document
                if not documents:
                    documents.append({
                        'text': f"PowerPoint file: {file_path.name} (no text extracted)",
                        'metadata': {
                            'source_file': Path(file_path).name,
                            'file_type': 'pptx',
                            'category': category,
                            'total_slides': len(prs.slides),
                            'empty_presentation': True
                        }
                    })

            except Exception as e:
                logger.warning(f"Could not extract text from PowerPoint: {e}")
                documents = [{
                    'text': f"PowerPoint file: {file_path.name} (text extraction failed: {str(e)})",
                    'metadata': {
                        'source_file': Path(file_path).name,
                        'file_type': 'pptx',
                        'category': category,
                        'file_size': os.path.getsize(file_path),
                        'extraction_failed': True
                    }
                }]

            # Store in PostgreSQL vector database
            connector = self.context_collections[category]
            success = connector.vector_store_documents(documents)

            if success:
                result['success'] = True
                result['document_count'] = len(documents)
            else:
                result['error'] = "Failed to store documents in vector database"

        except Exception as e:
            result['error'] = str(e)
            logger.error(f"Error processing PowerPoint file {file_path}: {e}")

        finally:
            result['processing_time'] = time.time() - start_time

        return result

    def _process_image_file(self, file_path, category):
        """Process image files and store in PostgreSQL"""
        start_time = time.time()
        result = {
            'success': False,
            'document_count': 0,
            'collections': [category],
            'context_type': category,
            'processing_time': 0,
            'error': None
        }

        try:
            # Parse the image using the LLM
            parsed_text = self.image_parser.parse_image(file_path)

            if parsed_text.startswith("Error:"):
                result['error'] = parsed_text
                logger.error(f"Error processing image file {file_path}: {parsed_text}")
                result['success'] = False
            else:
                document = {
                    'text': parsed_text,
                    'metadata': {
                        'source_file': file_path.name,
                        'file_type': 'image',
                        'category': category,
                        'file_size': os.path.getsize(file_path),
                    }
                }

                # Store in PostgreSQL vector database
                connector = self.context_collections[category]
                success = connector.vector_store_documents([document])

                if success:
                    result['success'] = True
                    result['document_count'] = 1
                else:
                    result['error'] = "Failed to store document in vector database"

        except Exception as e:
            result['error'] = str(e)
            logger.error(f"Error processing image file {file_path}: {e}")

        finally:
            result['processing_time'] = time.time() - start_time

        return result

    def _display_final_stats(self):
        """Display final processing statistics"""
        print(f"\n{Fore.GREEN}🎉 Context Processing Complete!{Style.RESET_ALL}")
        print(f"📊 Processing Statistics:")
        print(f"  • Total files found: {self.stats['total_files']}")
        print(f"  • Successfully processed: {Fore.GREEN}{self.stats['processed_files']}{Style.RESET_ALL}")
        print(f"  • Skipped files: {Fore.YELLOW}{self.stats['skipped_files']}{Style.RESET_ALL}")
        print(f"  • Failed files: {Fore.RED}{self.stats['failed_files']}{Style.RESET_ALL}")
        print(f"  • Total chunks created: {Fore.CYAN}{self.stats['total_chunks']}{Style.RESET_ALL}")
        print(f"  • Processing time: {Fore.CYAN}{self.stats['processing_time']:.2f}s{Style.RESET_ALL}")

        print(f"\n📊 By Category:")
        for category, stats in self.stats['by_category'].items():
            if stats['total'] > 0:
                print(
                    f"  • {category}: {stats['processed']} processed, {stats['skipped']} skipped, {stats['failed']} failed")

        # Display database statistics
        print(f"\n💾 Database Statistics:")
        try:
            for category, connector in self.context_collections.items():
                if category in self.stats['by_category'] and self.stats['by_category'][category]['processed'] > 0:
                    connector.diagnose_database()
        except Exception as e:
            print(f"{Fore.RED}Error getting database statistics: {e}{Style.RESET_ALL}")


def main():
    parser = argparse.ArgumentParser(
        description='Process context files for StorySense and store in PostgreSQL vector database'
    )

    parser.add_argument('--folder', '-f', required=False, default='../Input/ContextLibrary',
                        help='Path to context folder containing files to process')
    parser.add_argument('--force-reprocess', '-fr', action='store_true',
                        help='Force reprocessing of all files (ignore cache)')

    args = parser.parse_args()

    try:
        print(f"\n{Fore.CYAN}🚀 StorySense Context Processor with PostgreSQL Vector Storage{Style.RESET_ALL}")
        print(f"{'=' * 70}")

        # Validate folder path
        context_folder = Path(args.folder)
        print(f"Looking for context folder at: {context_folder.absolute()}")

        if not context_folder.exists():
            print(f"{Fore.RED}❌ Context folder not found: {context_folder}{Style.RESET_ALL}")
            print(f"{Fore.YELLOW}Creating directory structure...{Style.RESET_ALL}")

            # Create the directory
            os.makedirs(context_folder, exist_ok=True)

            # Create subdirectories for different context types
            subdirs = ['business_rules', 'requirements', 'documentation',
                       'policies', 'examples', 'glossary', 'wireframes']

            for subdir in subdirs:
                os.makedirs(context_folder / subdir, exist_ok=True)

            print(f"{Fore.GREEN}✅ Directory structure created at {context_folder.absolute()}{Style.RESET_ALL}")
            print(f"{Fore.YELLOW}Please add your context files and run again.{Style.RESET_ALL}")
            return

        # Initialize processor
        processor = EnhancedContextProcessor(context_folder)

        # Process context files
        stats = processor.process_all_context_files(force_reprocess=args.force_reprocess)

        # Display success message
        if stats['failed_files'] == 0:
            print(f"\n{Fore.GREEN}✅ All files processed successfully!{Style.RESET_ALL}")
        else:
            print(f"\n{Fore.YELLOW}⚠️  Processing completed with {stats['failed_files']} failures{Style.RESET_ALL}")

        print(
            f"{Fore.CYAN}📊 Processed {stats['processed_files']} files with {stats['total_chunks']} chunks{Style.RESET_ALL}")
        print(f"{Fore.CYAN}🎯 Ready for enhanced user story analysis!{Style.RESET_ALL}")

    except Exception as e:
        logger.error(f"Fatal error in context processor: {str(e)}")
        print(f"{Fore.RED}❌ Fatal error: {str(e)}{Style.RESET_ALL}")

        import traceback
        print(f"\n{Fore.CYAN}🐛 Debug Information:{Style.RESET_ALL}")
        print(f"{Fore.YELLOW}{traceback.format_exc()}{Style.RESET_ALL}")


if __name__ == "__main__":
    main()