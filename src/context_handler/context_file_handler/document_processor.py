import os
import json
import logging
from pathlib import Path
from typing import List, Dict, Any

# Document processing libraries
try:
    import PyPDF2
    import docx
    from pptx import Presentation
    import openpyxl
    import pandas as pd

    HAS_DOC_PROCESSORS = True
except Exception:
    HAS_DOC_PROCESSORS = False
    logging.warning(
        "Document processing libraries not installed. Install with: pip install PyPDF2 python-docx python-pptx openpyxl pandas")


class DocumentProcessor:
    """Parses many document types and returns a list of document dicts with text and metadata."""

    def __init__(self):
        self.file_processors = {
            '.pdf': self._process_pdf,
            '.docx': self._process_docx,
            '.doc': self._process_docx,
            '.pptx': self._process_pptx,
            '.xlsx': self._process_excel,
            '.xls': self._process_excel,
            '.txt': self._process_text,
            '.md': self._process_text,
            '.json': self._process_json,
            '.csv': self._process_csv
        }

    def process(self, file_path: str) -> List[Dict[str, Any]]:
        file_ext = Path(file_path).suffix.lower()

        if file_ext not in self.file_processors:
            raise ValueError(f"Unsupported file type: {file_ext}")

        processor = self.file_processors[file_ext]
        return processor(file_path)

    # --- processors ---
    def _process_pdf(self, file_path: str) -> List[Dict[str, Any]]:
        if not HAS_DOC_PROCESSORS:
            raise ImportError("PyPDF2 not installed")

        documents = []

        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)

            for page_num, page in enumerate(pdf_reader.pages):
                text = page.extract_text()
                if text and text.strip():
                    documents.append({
                        'text': text.strip(),
                        'metadata': {
                            'source_file': Path(file_path).name,
                            'file_type': 'pdf',
                            'page_number': page_num + 1,
                            'total_pages': len(pdf_reader.pages)
                        }
                    })

        return documents

    def _process_docx(self, file_path: str) -> List[Dict[str, Any]]:
        if not HAS_DOC_PROCESSORS:
            raise ImportError("python-docx not installed")

        documents = []
        doc = docx.Document(file_path)

        # Extract paragraphs
        full_text = []
        for paragraph in doc.paragraphs:
            if paragraph.text and paragraph.text.strip():
                full_text.append(paragraph.text.strip())

        # Extract tables
        table_content = []
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text and cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    table_content.append(" | ".join(row_text))

        # Combine content
        all_content = "\n".join(full_text)
        if table_content:
            all_content += "\n\nTables:\n" + "\n".join(table_content)

        if all_content.strip():
            documents.append({
                'text': all_content.strip(),
                'metadata': {
                    'source_file': Path(file_path).name,
                    'file_type': 'docx',
                    'paragraph_count': len(full_text),
                    'table_count': len(doc.tables)
                }
            })

        return documents

    def _process_pptx(self, file_path: str) -> List[Dict[str, Any]]:
        if not HAS_DOC_PROCESSORS:
            raise ImportError("python-pptx not installed")

        documents = []
        prs = Presentation(file_path)

        for slide_num, slide in enumerate(prs.slides):
            slide_text = []

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape.text.strip():
                    slide_text.append(shape.text.strip())

            if slide_text:
                documents.append({
                    'text': "\n".join(slide_text),
                    'metadata': {
                        'source_file': Path(file_path).name,
                        'file_type': 'pptx',
                        'slide_number': slide_num + 1,
                        'total_slides': len(prs.slides)
                    }
                })

        return documents

    def _process_excel(self, file_path: str) -> List[Dict[str, Any]]:
        if not HAS_DOC_PROCESSORS:
            raise ImportError("openpyxl not installed")

        documents = []

        try:
            import pandas as pd
            excel_file = pd.ExcelFile(file_path)

            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                text_content = df.to_string(index=False)

                if text_content.strip():
                    documents.append({
                        'text': text_content,
                        'metadata': {
                            'source_file': Path(file_path).name,
                            'file_type': 'xlsx',
                            'sheet_name': sheet_name,
                            'row_count': len(df),
                            'column_count': len(df.columns)
                        }
                    })

        except Exception as e:
            logging.error(f"Error processing Excel file {file_path}: {e}")
            raise

        return documents

    def _process_text(self, file_path: str) -> List[Dict[str, Any]]:
        documents = []

        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()

                if content and content.strip():
                    documents.append({
                        'text': content.strip(),
                        'metadata': {
                            'source_file': Path(file_path).name,
                            'file_type': Path(file_path).suffix.lower(),
                            'character_count': len(content),
                            'line_count': len(content.split('\n'))
                        }
                    })
        except UnicodeDecodeError:
            with open(file_path, 'r', encoding='latin-1') as file:
                content = file.read()

                if content and content.strip():
                    documents.append({
                        'text': content.strip(),
                        'metadata': {
                            'source_file': Path(file_path).name,
                            'file_type': Path(file_path).suffix.lower(),
                            'character_count': len(content),
                            'line_count': len(content.split('\n')),
                            'encoding': 'latin-1'
                        }
                    })

        return documents

    def _process_json(self, file_path: str) -> List[Dict[str, Any]]:
        documents = []

        with open(file_path, 'r', encoding='utf-8') as file:
            data = json.load(file)

            if isinstance(data, dict):
                text_content = self._dict_to_text(data)
            elif isinstance(data, list):
                text_content = "\n".join(
                    [self._dict_to_text(item) if isinstance(item, dict) else str(item) for item in data])
            else:
                text_content = str(data)

            if text_content.strip():
                documents.append({
                    'text': text_content.strip(),
                    'metadata': {
                        'source_file': Path(file_path).name,
                        'file_type': 'json',
                        'data_type': type(data).__name__
                    }
                })

        return documents

    def _process_csv(self, file_path: str) -> List[Dict[str, Any]]:
        documents = []

        import pandas as pd
        df = pd.read_csv(file_path)

        text_content = df.to_string(index=False)

        if text_content.strip():
            documents.append({
                'text': text_content,
                'metadata': {
                    'source_file': Path(file_path).name,
                    'file_type': 'csv',
                    'row_count': len(df),
                    'column_count': len(df.columns),
                    'columns': list(df.columns)
                }
            })

        return documents

    def _dict_to_text(self, data: dict, prefix: str = "") -> str:
        text_parts = []

        for key, value in data.items():
            if isinstance(value, dict):
                text_parts.append(f"{prefix}{key}:")
                text_parts.append(self._dict_to_text(value, prefix + "  "))
            elif isinstance(value, list):
                text_parts.append(f"{prefix}{key}: {', '.join(str(v) for v in value)}")
            else:
                text_parts.append(f"{prefix}{key}: {value}")

        return "\n".join(text_parts)
