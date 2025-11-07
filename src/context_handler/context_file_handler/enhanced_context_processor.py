import os
import logging
import time
import json
from pathlib import Path
from colorama import Fore, Style
import pandas as pd
import numpy as np
from datetime import datetime
import hashlib
import mimetypes
import boto3
import configparser
from src.context_handler.context_storage_handler.pgvector_connector import PGVectorConnector
from src.aws_layer.aws_titan_embedding import AWSTitanEmbeddings
from src.metrics.metrics_manager  import MetricsManager
# from src.configuration_handler.env_manager  import EnvManager

# Set up logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)


class EnhancedContextProcessor:
    """Enhanced context processor for multi-format document processing"""

    def __init__(self, context_folder_path, metrics_manager=None):
        self.context_folder_path = Path(context_folder_path)
        # self.env_manager = EnvManager()
        self.metrics_manager = metrics_manager or MetricsManager()

        # Initialize PGVector connectors for different context types
        self.context_collections = {
            'business_rules': PGVectorConnector(collection_name="business_rules", metrics_manager=self.metrics_manager),
            'requirements': PGVectorConnector(collection_name="requirements", metrics_manager=self.metrics_manager),
            'documentation': PGVectorConnector(collection_name="documentation", metrics_manager=self.metrics_manager),
            'policies': PGVectorConnector(collection_name="policies", metrics_manager=self.metrics_manager),
            'examples': PGVectorConnector(collection_name="examples", metrics_manager=self.metrics_manager),
            'glossary': PGVectorConnector(collection_name="glossary", metrics_manager=self.metrics_manager),
            'wireframes': PGVectorConnector(collection_name="wireframes", metrics_manager=self.metrics_manager)
        }

        # Initialize AWS Titan embeddings
        self.embeddings = AWSTitanEmbeddings(
            model_id=os.getenv('EMBEDDING_MODEL_NAME', 'amazon.titan-embed-text-v1'),
            local_storage_path=os.getenv('LOCAL_EMBEDDINGS_PATH', './Data/LocalEmbeddings')
        )

        # Supported formats
        self.text_formats = ['.xlsx', '.xls', '.csv', '.docx', '.doc', '.pdf', '.json', '.xml', '.txt', '.md', '.pptx',
                             '.ppt']
        self.image_formats = ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp']

        # Processing statistics
        self.stats = {
            'total_files': 0,
            'processed_files': 0,
            'skipped_files': 0,
            'failed_files': 0,
            'total_chunks': 0,
            'processing_time': 0,
            'by_category': {}
        }

        # Document registry to track processed files
        self.registry_file = Path('../Data/DocumentRegistry/document_registry.json')
        self.document_registry = self._load_document_registry()

        # Create necessary directories
        os.makedirs('../Data/DocumentRegistry', exist_ok=True)
        os.makedirs('../Data/LocalEmbeddings', exist_ok=True)

    def _load_document_registry(self):
        """Load the document registry that tracks processed files"""
        if self.registry_file.exists():
            try:
                with open(self.registry_file, 'r', encoding='utf-8') as f:
                    return json.load(f)
            except (json.JSONDecodeError, Exception) as e:
                logger.warning(f"Could not load document registry: {e}. Creating new registry.")
                return {}
        return {}

    def _save_document_registry(self):
        """Save the document registry"""
        try:
            with open(self.registry_file, 'w', encoding='utf-8') as f:
                json.dump(self.document_registry, f, indent=2, default=str)
        except Exception as e:
            logger.error(f"Could not save document registry: {e}")

    def _calculate_file_hash(self, file_path):
        """Calculate MD5 hash of a file for duplicate detection"""
        hash_md5 = hashlib.md5()
        try:
            with open(file_path, "rb") as f:
                for chunk in iter(lambda: f.read(4096), b""):
                    hash_md5.update(chunk)
            return hash_md5.hexdigest()
        except Exception as e:
            logger.error(f"Could not calculate hash for {file_path}: {e}")
            return ""

    def _is_document_processed(self, file_path):
        """Check if a document has already been processed"""
        file_hash = self._calculate_file_hash(file_path)
        file_name = Path(file_path).name
        file_size = os.path.getsize(file_path)
        file_mtime = os.path.getmtime(file_path)

        # Check by hash first (most reliable)
        for doc_id, doc_info in self.document_registry.items():
            if doc_info.get('file_hash') == file_hash and file_hash:
                return {
                    'exists': True,
                    'reason': 'identical_content',
                    'existing_doc_id': doc_id,
                    'existing_info': doc_info
                }

        # Check by name, size, and modification time
        for doc_id, doc_info in self.document_registry.items():
            if (doc_info.get('file_name') == file_name and
                    doc_info.get('file_size') == file_size and
                    doc_info.get('file_mtime') == file_mtime):
                return {
                    'exists': True,
                    'reason': 'same_file_attributes',
                    'existing_doc_id': doc_id,
                    'existing_info': doc_info
                }

        return {'exists': False}

    def _register_document(self, file_path, processing_result):
        """Register a processed document in the registry"""
        doc_id = f"doc_{datetime.now().strftime('%Y%m%d_%H%M%S')}_{len(self.document_registry)}"

        self.document_registry[doc_id] = {
            'file_name': Path(file_path).name,
            'file_path': str(file_path),
            'file_hash': self._calculate_file_hash(file_path),
            'file_size': os.path.getsize(file_path),
            'file_mtime': os.path.getmtime(file_path),
            'file_type': Path(file_path).suffix.lower(),
            'processed_date': datetime.now().isoformat(),
            'document_count': processing_result.get('document_count', 0),
            'collections': processing_result.get('collections', []),
            'context_type': processing_result.get('context_type', 'unknown'),
            'processing_time': processing_result.get('processing_time', 0),
            'success': processing_result.get('success', False)
        }

        self._save_document_registry()
        return doc_id

    def process_all_context_files(self, force_reprocess=False):
        """Process all context files in the folder and subfolders"""
        print(f"\n{Fore.CYAN}🚀 Starting Enhanced Context Processing{Style.RESET_ALL}")
        print(f"📁 Context Folder: {self.context_folder_path}")

        start_time = time.time()

        if not self.context_folder_path.exists():
            raise FileNotFoundError(f"Context folder not found: {self.context_folder_path}")

        # Discover all files
        all_files = self._discover_files()
        self.stats['total_files'] = sum(len(files) for files in all_files.values())

        print(f"📊 Found {self.stats['total_files']} files to process")

        # Process files by category
        for category, files in all_files.items():
            if files:
                print(f"\n{Fore.YELLOW}📂 Processing {category} files ({len(files)} files){Style.RESET_ALL}")
                self._process_category_files(category, files, force_reprocess)

        # Calculate final statistics
        self.stats['processing_time'] = time.time() - start_time
        self._display_final_stats()

        return self.stats

    def _discover_files(self):
        """Discover and categorize all files in the context folder"""
        categorized_files = {
            'business_rules': [],
            'requirements': [],
            'documentation': [],
            'policies': [],
            'examples': [],
            'glossary': [],
            'wireframes': []
        }

        for root, dirs, files in os.walk(self.context_folder_path):
            category_hint = Path(root).name.lower()

            for file in files:
                file_path = Path(root) / file
                file_ext = file_path.suffix.lower()

                # Skip hidden files and system files
                if file.startswith('.') or file.startswith('~'):
                    continue

                # Skip files that are too large
                max_size_mb = float(os.getenv('MAX_FILE_SIZE_MB', '500'))
                if os.path.getsize(file_path) > max_size_mb * 1024 * 1024:
                    logger.warning(f"Skipping large file: {file_path} (exceeds {max_size_mb}MB)")
                    continue

                # Categorize based on extension and folder name
                category = self._categorize_file(file_path, category_hint)
                categorized_files[category].append(file_path)

        return categorized_files

    def _categorize_file(self, file_path, folder_hint):
        """Categorize file based on extension and folder context"""
        file_ext = file_path.suffix.lower()
        file_name = file_path.name.lower()

        # Check folder name first for strongest hint
        if any(hint in folder_hint for hint in ['business', 'rule', 'logic', 'constraint']):
            return 'business_rules'
        elif any(hint in folder_hint for hint in ['requirement', 'spec', 'functional', 'non-functional']):
            return 'requirements'
        elif any(hint in folder_hint for hint in ['doc', 'guide', 'manual', 'readme']):
            return 'documentation'
        elif any(hint in folder_hint for hint in ['policy', 'procedure', 'guideline', 'standard']):
            return 'policies'
        elif any(hint in folder_hint for hint in ['example', 'template', 'sample', 'demo']):
            return 'examples'
        elif any(hint in folder_hint for hint in ['glossary', 'term', 'definition', 'acronym']):
            return 'glossary'
        elif any(hint in folder_hint for hint in ['wireframe', 'mockup', 'ui', 'design']):
            return 'wireframes'

        # Then check file name
        if any(keyword in file_name for keyword in ['business', 'rule', 'logic', 'constraint']):
            return 'business_rules'
        elif any(keyword in file_name for keyword in ['requirement', 'spec', 'functional', 'non-functional']):
            return 'requirements'
        elif any(keyword in file_name for keyword in ['doc', 'guide', 'manual', 'readme']):
            return 'documentation'
        elif any(keyword in file_name for keyword in ['policy', 'procedure', 'guideline', 'standard']):
            return 'policies'
        elif any(keyword in file_name for keyword in ['example', 'template', 'sample', 'demo']):
            return 'examples'
        elif any(keyword in file_name for keyword in ['glossary', 'term', 'definition', 'acronym']):
            return 'glossary'
        elif any(keyword in file_name for keyword in ['wireframe', 'mockup', 'ui', 'design']):
            return 'wireframes'

        # Default based on file extension
        if file_ext in self.image_formats:
            return 'wireframes' if 'ui' in file_name or 'screen' in file_name else 'documentation'
        elif file_ext in ['.xlsx', '.xls', '.csv']:
            return 'glossary' if 'glossary' in file_name else 'business_rules'
        elif file_ext in ['.pdf', '.docx', '.doc']:
            return 'documentation'
        elif file_ext in ['.pptx', '.ppt']:
            return 'requirements'
        else:
            return 'documentation'  # Default category

    def _process_category_files(self, category, files, force_reprocess):
        """Process files in a specific category"""
        if category not in self.stats['by_category']:
            self.stats['by_category'][category] = {
                'total': len(files),
                'processed': 0,
                'skipped': 0,
                'failed': 0
            }

        for file_path in files:
            try:
                print(f"  📄 Processing: {file_path.name}")

                # Check if file already processed
                if not force_reprocess:
                    duplicate_check = self._is_document_processed(file_path)
                    if duplicate_check['exists']:
                        print(f"    ⏭️  Skipping {Path(file_path).name} - {duplicate_check['reason']}")
                        self.stats['skipped_files'] += 1
                        self.stats['by_category'][category]['skipped'] += 1
                        continue

                # Process based on file extension
                file_ext = file_path.suffix.lower()

                if file_ext in ['.xlsx', '.xls', '.csv']:
                    result = self._process_structured_document(file_path, category)
                elif file_ext in ['.pdf', '.docx', '.doc', '.txt', '.md']:
                    result = self._process_unstructured_document(file_path, category)
                elif file_ext in ['.pptx', '.ppt']:
                    result = self._process_presentation(file_path, category)
                elif file_ext in self.image_formats:
                    result = self._process_image_file(file_path, category)
                else:
                    print(f"    ⚠️  Skipping unsupported file type: {file_ext}")
                    self.stats['skipped_files'] += 1
                    self.stats['by_category'][category]['skipped'] += 1
                    continue

                if result['success']:
                    # Register the document
                    doc_id = self._register_document(file_path, result)

                    self.stats['processed_files'] += 1
                    self.stats['by_category'][category]['processed'] += 1
                    self.stats['total_chunks'] += result['document_count']

                    print(f"    ✅ Successfully processed {result['document_count']} chunks")
                else:
                    self.stats['failed_files'] += 1
                    self.stats['by_category'][category]['failed'] += 1
                    print(f"    ❌ Failed: {result.get('error', 'Unknown error')}")

            except Exception as e:
                logger.error(f"Error processing {file_path}: {str(e)}")
                print(f"    ❌ Failed: {str(e)}")
                self.stats['failed_files'] += 1
                self.stats['by_category'][category]['failed'] += 1

    def _process_structured_document(self, file_path, category):
        """Process structured documents (Excel, CSV)"""
        start_time = time.time()
        result = {
            'success': False,
            'document_count': 0,
            'collections': [category],
            'context_type': category,
            'processing_time': 0,
            'error': None
        }

        try:
            # Read the file based on extension
            file_ext = file_path.suffix.lower()
            if file_ext in ['.xlsx', '.xls']:
                df = pd.read_excel(file_path)
            elif file_ext == '.csv':
                df = pd.read_csv(file_path)
            else:
                raise ValueError(f"Unsupported structured document format: {file_ext}")

            # Convert to documents
            documents = []

            # Process as table
            if len(df) > 0:
                # Process each row as a separate document
                for idx, row in df.iterrows():
                    content = self._row_to_text(row)
                    if content.strip():
                        documents.append({
                            'text': content,
                            'metadata': {
                                'source_file': file_path.name,
                                'file_type': 'structured_document',
                                'row_index': idx
                            }
                        })

            # Store documents in vector database
            if documents:
                connector = self.context_collections[category]
                success = connector.vector_store_documents(documents)

                if success:
                    result['success'] = True
                    result['document_count'] = len(documents)
                else:
                    result['error'] = "Failed to store documents in vector database"
            else:
                result['error'] = "No content extracted from file"

        except Exception as e:
            result['error'] = str(e)
            logger.error(f"Error processing structured document {file_path}: {e}")

        finally:
            result['processing_time'] = time.time() - start_time

        return result

    def _process_unstructured_document(self, file_path, category):
        """Process unstructured documents (PDF, DOCX, TXT)"""
        start_time = time.time()
        result = {
            'success': False,
            'document_count': 0,
            'collections': [category],
            'context_type': category,
            'processing_time': 0,
            'error': None
        }

        try:
            file_ext = file_path.suffix.lower()

            if file_ext == '.pdf':
                documents = self._process_pdf(file_path)
            elif file_ext in ['.docx', '.doc']:
                documents = self._process_docx(file_path)
            elif file_ext in ['.txt', '.md']:
                documents = self._process_text(file_path)
            else:
                raise ValueError(f"Unsupported unstructured document format: {file_ext}")

            # Store documents in vector database
            if documents:
                connector = self.context_collections[category]
                success = connector.vector_store_documents(documents)

                if success:
                    result['success'] = True
                    result['document_count'] = len(documents)
                else:
                    result['error'] = "Failed to store documents in vector database"
            else:
                result['error'] = "No content extracted from file"

        except Exception as e:
            result['error'] = str(e)
            logger.error(f"Error processing unstructured document {file_path}: {e}")

        finally:
            result['processing_time'] = time.time() - start_time

        return result

    def _process_presentation(self, file_path, category):
        """Process presentation files (PPTX, PPT)"""
        start_time = time.time()
        result = {
            'success': False,
            'document_count': 0,
            'collections': [category],
            'context_type': category,
            'processing_time': 0,
            'error': None
        }

        try:
            documents = self._process_pptx(file_path)

            # Store documents in vector database
            if documents:
                connector = self.context_collections[category]
                success = connector.vector_store_documents(documents)

                if success:
                    result['success'] = True
                    result['document_count'] = len(documents)
                else:
                    result['error'] = "Failed to store documents in vector database"
            else:
                result['error'] = "No content extracted from file"

        except Exception as e:
            result['error'] = str(e)
            logger.error(f"Error processing presentation {file_path}: {e}")

        finally:
            result['processing_time'] = time.time() - start_time

        return result

    def _process_image_file(self, file_path, category):
        """Process image files with OCR and metadata extraction"""
        start_time = time.time()
        result = {
            'success': False,
            'document_count': 0,
            'collections': [category],
            'context_type': category,
            'processing_time': 0,
            'error': None
        }

        try:
            # Extract metadata
            metadata = self._extract_image_metadata(file_path)

            # OCR text extraction (if enabled)
            ocr_text = ""
            if self.env_manager.get_env('OCR_ENABLED', 'false').lower() == 'true':
                ocr_text = self._extract_text_from_image(file_path)

            # Create content
            content = f"Image: {file_path.name}\n"
            content += f"Type: {metadata.get('format', 'Unknown')}\n"
            content += f"Size: {metadata.get('size', 'Unknown')}\n"
            if ocr_text:
                content += f"Extracted Text: {ocr_text}\n"

            documents = [{
                'text': content,
                'metadata': {
                    'source_file': file_path.name,
                    'file_type': 'image',
                    'image_metadata': metadata
                }
            }]

            # Store documents in vector database
            if documents:
                connector = self.context_collections[category]
                success = connector.vector_store_documents(documents)

                if success:
                    result['success'] = True
                    result['document_count'] = len(documents)
                else:
                    result['error'] = "Failed to store documents in vector database"
            else:
                result['error'] = "No content extracted from file"

        except Exception as e:
            result['error'] = str(e)
            logger.error(f"Error processing image {file_path}: {e}")

        finally:
            result['processing_time'] = time.time() - start_time

        return result

    def _row_to_text(self, row):
        """Convert DataFrame row to text content"""
        content_parts = []
        for col, val in row.items():
            if pd.notna(val) and str(val).strip():
                content_parts.append(f"{col}: {val}")
        return "\n".join(content_parts)

    def _extract_image_metadata(self, file_path):
        """Extract metadata from image files"""
        try:
            from PIL import Image
            with Image.open(file_path) as img:
                return {
                    'format': img.format,
                    'mode': img.mode,
                    'size': f"{img.width}x{img.height}",
                    'file_size': os.path.getsize(file_path)
                }
        except Exception:
            return {'format': 'Unknown', 'size': 'Unknown'}

    def _extract_text_from_image(self, file_path):
        """Extract text from image using OCR"""
        try:
            import pytesseract
            from PIL import Image

            with Image.open(file_path) as img:
                text = pytesseract.image_to_string(img)
                return text.strip()
        except Exception as e:
            logger.warning(f"OCR failed for {file_path}: {str(e)}")
            return ""

    def _process_pdf(self, file_path):
        """Process PDF files"""
        try:
            import PyPDF2
            documents = []

            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)

                for page_num, page in enumerate(pdf_reader.pages):
                    text = page.extract_text()
                    if text.strip():
                        documents.append({
                            'text': text.strip(),
                            'metadata': {
                                'source_file': Path(file_path).name,
                                'file_type': 'pdf',
                                'page_number': page_num + 1,
                                'total_pages': len(pdf_reader.pages)
                            }
                        })

            return documents
        except ImportError:
            logger.warning("PyPDF2 not installed. Install with: pip install PyPDF2")
            return []
        except Exception as e:
            logger.error(f"Error processing PDF {file_path}: {e}")
            return []

    def _process_docx(self, file_path):
        """Process Word documents"""
        try:
            import docx
            documents = []
            doc = docx.Document(file_path)

            # Extract paragraphs
            full_text = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    full_text.append(paragraph.text.strip())

            # Extract tables
            table_content = []
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        table_content.append(" | ".join(row_text))

            # Combine content
            all_content = "\n".join(full_text)
            if table_content:
                all_content += "\n\nTables:\n" + "\n".join(table_content)

            if all_content.strip():
                documents.append({
                    'text': all_content.strip(),
                    'metadata': {
                        'source_file': Path(file_path).name,
                        'file_type': 'docx',
                        'paragraph_count': len(full_text),
                        'table_count': len(doc.tables)
                    }
                })

            return documents
        except ImportError:
            logger.warning("python-docx not installed. Install with: pip install python-docx")
            return []
        except Exception as e:
            logger.error(f"Error processing DOCX {file_path}: {e}")
            return []

    def _process_pptx(self, file_path):
        """Process PowerPoint presentations"""
        try:
            from pptx import Presentation
            documents = []
            prs = Presentation(file_path)

            for slide_num, slide in enumerate(prs.slides):
                slide_text = []

                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())

                if slide_text:
                    documents.append({
                        'text': "\n".join(slide_text),
                        'metadata': {
                            'source_file': Path(file_path).name,
                            'file_type': 'pptx',
                            'slide_number': slide_num + 1,
                            'total_slides': len(prs.slides)
                        }
                    })

            return documents
        except ImportError:
            logger.warning("python-pptx not installed. Install with: pip install python-pptx")
            return []
        except Exception as e:
            logger.error(f"Error processing PPTX {file_path}: {e}")
            return []

    def _process_text(self, file_path):
        """Process text and markdown files"""
        documents = []

        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()

                if content.strip():
                    documents.append({
                        'text': content.strip(),
                        'metadata': {
                            'source_file': Path(file_path).name,
                            'file_type': Path(file_path).suffix.lower(),
                            'character_count': len(content),
                            'line_count': len(content.split('\n'))
                        }
                    })
        except UnicodeDecodeError:
            # Try with different encoding
            with open(file_path, 'r', encoding='latin-1') as file:
                content = file.read()

                if content.strip():
                    documents.append({
                        'text': content.strip(),
                        'metadata': {
                            'source_file': Path(file_path).name,
                            'file_type': Path(file_path).suffix.lower(),
                            'character_count': len(content),
                            'line_count': len(content.split('\n')),
                            'encoding': 'latin-1'
                        }
                    })

        return documents

    def _display_final_stats(self):
        """Display final processing statistics"""
        print(f"\n{Fore.GREEN}🎉 Context Processing Complete!{Style.RESET_ALL}")
        print(f"📊 Processing Statistics:")
        print(f"  • Total files found: {self.stats['total_files']}")
        print(f"  • Successfully processed: {Fore.GREEN}{self.stats['processed_files']}{Style.RESET_ALL}")
        print(f"  • Skipped files: {Fore.YELLOW}{self.stats['skipped_files']}{Style.RESET_ALL}")
        print(f"  • Failed files: {Fore.RED}{self.stats['failed_files']}{Style.RESET_ALL}")
        print(f"  • Total chunks created: {Fore.CYAN}{self.stats['total_chunks']}{Style.RESET_ALL}")
        print(f"  • Processing time: {Fore.CYAN}{self.stats['processing_time']:.2f}s{Style.RESET_ALL}")

        print(f"\n📊 By Category:")
        for category, stats in self.stats['by_category'].items():
            if stats['total'] > 0:
                print(
                    f"  • {category}: {stats['processed']} processed, {stats['skipped']} skipped, {stats['failed']} failed")